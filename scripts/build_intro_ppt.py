"""Generate English introduction deck with diagrams and project-generated examples."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "doc" / "ppt_assets"
OUT = ROOT / "doc" / "PowerSeqGen_Introduction_v1.5.pptx"

# Theme
BG = RGBColor(0x0F, 0x17, 0x2A)
TITLE = RGBColor(0xF8, 0xFA, 0xFC)
BODY = RGBColor(0xCB, 0xD5, 0xE1)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def _fill_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.95))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TITLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(4)


def _add_bullets(slide, items: list[str], *, left=0.6, top=1.35, width=5.8, height=5.6, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith("  ") else 0
        p.font.size = Pt(size - 2 if p.level else size)
        p.font.color.rgb = MUTED if p.level else BODY
        p.space_after = Pt(6)


def _add_image(slide, path: Path, *, left, top, width, height=None):
    if not path.is_file():
        return
    if height is not None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def _add_image_fit(slide, path: Path, *, left, top, width, height):
    if not path or not path.is_file():
        return
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_w, box_h = width, height
    scale = min(box_w / img_w, box_h / img_h)
    final_w = img_w * scale
    final_h = img_h * scale
    final_left = left + (box_w - final_w) / 2
    final_top = top + (box_h - final_h) / 2
    slide.shapes.add_picture(
        str(path),
        Inches(final_left),
        Inches(final_top),
        width=Inches(final_w),
        height=Inches(final_h),
    )


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    return slide


def _title_slide(prs: Presentation, pipeline: Path | None):
    slide = _blank(prs)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.8), Inches(2.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Power Sequence Generator"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE
    p2 = tf.add_paragraph()
    p2.text = "From dependency graph to RTL, firmware, and documentation"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "pwrseq_gen  ·  v1.5  ·  PSEQCELL-based power sequencing"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(20)
    if pipeline and pipeline.is_file():
        _add_image(slide, pipeline, left=1.2, top=4.0, width=10.8)


def _image_slide(prs, title, subtitle, image: Path | None, caption: str | None = None):
    slide = _blank(prs)
    _add_title(slide, title, subtitle)
    if image and image.is_file():
        _add_image_fit(slide, image, left=0.55, top=1.25, width=12.2, height=5.5)
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(11)
        cp.font.color.rgb = MUTED
        cp.alignment = PP_ALIGN.CENTER


def _split_slide(prs, title, subtitle, bullets, image: Path | None, *, img_left=6.35, img_width=6.5):
    slide = _blank(prs)
    _add_title(slide, title, subtitle)
    _add_bullets(slide, bullets, width=5.45)
    if image and image.is_file():
        _add_image(slide, image, left=img_left, top=1.25, width=img_width, height=5.55)


def build(assets: dict[str, Path] | None = None) -> Path:
    if assets is None:
        sys.path.insert(0, str(ROOT / "scripts"))
        from ppt_asset_gen import generate_all_assets

        assets = generate_all_assets()

    def a(key: str) -> Path | None:
        p = assets.get(key)
        return p if p and p.is_file() else None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs, a("pipeline"))

    _split_slide(
        prs,
        "Why this tool?",
        "One model — many deliverables",
        [
            "Power-up/down has many rails, PG signals, and timing constraints.",
            "Specs drift when JSON, RTL, firmware, and diagrams are edited by hand.",
            "pwrseq_gen keeps one JSON source of truth, such as doc/demo_json.json.",
            "  It generates Verilog, C, Draw.io, and WaveDrom from the same graph.",
        ],
        a("pipeline"),
    )

    _image_slide(
        prs,
        "End-to-end pipeline",
        "Author in GUI or JSON → validate → generate / export",
        a("pipeline"),
        "The provided doc/demo_* files show the resulting JSON, Verilog, C, Draw.io, and WaveDrom artifacts.",
    )

    _split_slide(
        prs,
        "Configuration model",
        "Rails, conditions, validation",
        [
            "Node = Input (debounce) or Output (PSEQCELL).",
            "Hi / Lo / Force: AND inside group, OR across groups.",
            "Invert & use-mode (node / Hi / Lo / self feedback).",
            "Timing: CYCLE_*, RECOVER, INIT, FORCE, OD.",
            "Validator catches duplicates, missing refs, cycles.",
        ],
        a("config_model"),
        img_left=5.9,
        img_width=6.9,
    )

    _image_slide(
        prs,
        "Provided input spec",
        "doc/demo_json.json — project-level power sequence configuration",
        a("demo_json_snippet"),
        "This user-provided demo config is the primary example used in this deck.",
    )

    _image_slide(
        prs,
        "Provided Draw.io output",
        "doc/demo_drawio_png.png — dependency graph rendered from the demo config",
        a("demo_drawio_png"),
        "This is the provided Draw.io diagram asset, not a replacement generated from output/test.json.",
    )

    _image_slide(
        prs,
        "Draw.io routing details",
        "doc/demo_drawio_png.png — layered placement, Hi/Lo color paths, and feedback routing",
        a("demo_drawio_zoom") or a("demo_drawio_png"),
        "Zoomed from the provided Draw.io PNG: Input → logic gates → PSEQCELL blocks → output labels.",
    )

    _image_slide(
        prs,
        "Draw.io layer architecture",
        "Left → right signal flow (doc/DRAWIO_RULES.md)",
        a("layers"),
    )

    _split_slide(
        prs,
        "PSEQCELL in RTL",
        "Aligned with src/reference/PSEQCELL.v",
        [
            "Ports: iHi, iLo, iForce → internal seq FSM.",
            "o maps to oRAIL; ~Q used for feedback paths.",
            "Parameters from JSON: CYCLE_HI/LO, RECOVER, etc.",
            "DEB instances for debounced inputs.",
        ],
        a("pseqcell"),
        img_left=6.2,
        img_width=6.5,
    )

    _image_slide(
        prs,
        "Verilog generation",
        "Auto-generated module + PSEQCELL instances",
        a("demo_verilog_snippet"),
        "Source asset: doc/demo_verilog.v. The generator derives ports, condition wires, DEB, and PSEQCELL instances from JSON.",
    )

    _image_slide(
        prs,
        "C firmware generation",
        "pwrcell_t scaffolding (reference: pwrcell.c/h)",
        a("demo_c_snippet"),
        "Source asset: doc/demo_c.c. The output contains per-rail pwrcell_t state and runtime condition hooks.",
    )

    _image_slide(
        prs,
        "Provided WaveDrom timing",
        "Simulated from dependency graph (not RTL cycle-accurate)",
        a("demo_wavedrom_png"),
        "Source asset: doc/demo_wavedrom_png.png, with JSON data in doc/demo_wavedrow_json.json.",
    )

    _image_slide(
        prs,
        "WaveDrom review use case",
        "doc/demo_wavedrom_png.png — 50-step sequencing overview",
        a("demo_wavedrom_zoom") or a("demo_wavedrom_png"),
        "Zoomed from the provided WaveDrom PNG for dependency arrows and power-up timing review.",
    )

    _image_slide(
        prs,
        "Tooling & quality",
        "Build, test, package, and review",
        a("demo_drawio_png"),
        "Python 3.8+, CustomTkinter GUI, pytest coverage, PyInstaller packaging, and reference libraries.",
    )

    slide = _blank(prs)
    _add_title(slide, "Getting started", "Five steps")
    steps = [
        "1. pip install -r requirements.txt  (or run run.bat)",
        "2. python src/main.py — build nodes & conditions",
        "3. Confirm green validation bar",
        "4. Generate → Verilog / C  |  Export → Draw.io / WaveDrom",
        "5. json2drawio.bat doc\\demo_json.json -o doc\\demo_drawio_xml.xml",
    ]
    _add_bullets(slide, steps, width=12.0, top=1.4)
    if a("pipeline"):
        _add_image(slide, a("pipeline"), left=0.9, top=4.45, width=11.5)

    slide = _blank(prs)
    _add_title(slide, "Summary", None)
    if a("pipeline"):
        _add_image(slide, a("pipeline"), left=0.9, top=1.35, width=11.6)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(5.25), Inches(11.2), Inches(1.25))
    tf = box.text_frame
    for i, line in enumerate(
        [
            "One JSON model captures power-sequence intent",
            "Code, diagrams, and timing documentation stay synchronized",
            "",
            "Questions?",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if i < 2 else 28)
        p.font.color.rgb = ACCENT if i == 3 else BODY
        p.font.bold = i == 3
        p.space_after = Pt(10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUT))
        return OUT
    except PermissionError:
        alt = OUT.with_name(OUT.stem + "_rebuilt.pptx")
        prs.save(str(alt))
        print(f"Warning: {OUT.name} is locked — wrote {alt.name} instead", file=sys.stderr)
        return alt


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
    sys.exit(0)
